#Requirements:
#project must have a main funciton and three or more additional functions.                                  done
#atleast 3 of the additional functions must be accompanied by tests with pytest.                            done
#test functions must be in test_project.py.                                                                 done
#project should take more time than the problem sets                                                        done
#any pip installs should be listed in a document called requirements.txt in the root of the project.        done

'''
The idea

automatic powerpoint generator with pictures, productspecifications and the titles from products sold on Amazon.de
reads from a list of product urls and outputs them into a powerpoint presentation.
''' 

#Plan:
#open links.txt file and read the links from it.
#use re.search on them, to make sure they are valid links
#if one link doesn't work sys.exit with the link that doesn't work ONLY IF ALL LINKS WORK SHOULD IT GENERATE A PPTX FILE!!
#scrape function to scrape information from given urls
#pptx function to paste what has been scraped and output a pp
#ask user for name for pp file


import pptx
from pptx.util import Inches, Cm, Pt
import sys
import re
import requests
import time
from bs4 import BeautifulSoup
from pptx.dml.color import RGBColor
import os
from headers import headers as headers1

def main():
    #prompts user for filename and runs all other functions
    filename = input("Filename: ").strip()

    pptitle = "Source Direct,\nBuy-Direct"
    ppsubtitle = "Your sourcing office in Asia"
    bottomrighttitle = pptitle
    maincolor = RGBColor(51, 51, 102) #change RGB code to change color
    secondcolor = RGBColor(248, 98, 57) #change RGB code to change color
    #If you want to use your own logo, paste your logo into the rootfolder of this script and call it: logo.png
    #Make sure to save before closing

    urls = linkreader()
    linkchecker(urls)
    information, productnumber = scraper(urls)
    presentationator(filename, information, pptitle, ppsubtitle, bottomrighttitle, maincolor, secondcolor)
    removed = imageremover(productnumber)
    if removed:
        print("All images have been removed :)")
    if not removed:
        print("Not all images were removed :(")
    print("All done!")


def linkreader():
    #reads all urls in linkts.txt
    with open("links.txt", "r") as file:
        urls = file.readlines()
        return urls


def linkchecker(urls):
    #checks validity of urls
    line = 1
    for url in urls:
        if re.search(r"^(?:https?://)?(?:www\.)?(amazon\.de/).*", url, re.IGNORECASE):
            line += 1
            continue
        else:
            sys.exit(f"Sorry... {url}, on line {line} is not a valid amazon.de link.")
    return urls


def scraper(urls):
    #scrapes all urls for titles, pictures, prices and descriptions
    information = []
    #for counting
    productnumber = 0
    for url in urls:
        productnumber += 1

        #beautifulsoup
        headers = headers1
        result = requests.get(url, headers=headers)
        doc = BeautifulSoup(result.text, "html.parser")

        #title, price, url
        title = doc.find(id="productTitle").get_text()
        if short := re.search(r"^([a-zA-Z0-9À-ÿ -]+),?\.?\|?(?: - )?", title.strip()):
            short_title = short.group(1).strip()
        price = doc.find("span", {"class": "a-offscreen"}).get_text()
        link = url
        if "\n" in url:
            link = link.replace("\n","")

        #image
        image = doc.select_one("#landingImage")
        image = image.attrs.get("src")
        with open("AmazonProductImage"+str(productnumber)+".jpg", "wb") as file:
            im = requests.get(image)
            file.write(im.content)
            print(f"Image {productnumber} scraped!")

        #characteristics
        index = 0
        characteristics = doc.find("table", {"class":"a-normal a-spacing-micro"}).get_text()#finds table with characteristics
        characteristics = characteristics.strip().replace("   ","_").replace("     ","_").split("_")
        char_list = [] 
        for c in characteristics[:-1]:
            if index == 0 or index%2 == 0:
                char_list.append({c.strip():characteristics[index+1].strip()})
            index += 1
        
        information.append({"title":short_title, "price":price, "url":link, "characteristics": char_list})
        time.sleep(2) #2 second wait period, so program doesn't overload Amazon. Amazon blocks bots.

    return information, productnumber


def presentationator(filename, information, pptitle, ppsubtitle, bottomrighttitle, maincolor, secondcolor):
    #puts all scraped information into powerpoint slides
    #title slide
    pres = pptx.Presentation()
    title_slide_layout = pres.slide_layouts[0] #Title slide layout
    slide = pres.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = pptitle
    subtitle.text = ppsubtitle #subtitle text
    #color codes should be rgb(51, 51, 102) and rgb(248, 98, 57)
    subtitle.text_frame.paragraphs[0].font.color.rgb = maincolor
    title_para = slide.shapes.title.text_frame.paragraphs[0]
    title_para1 = slide.shapes.title.text_frame.paragraphs[1]
    title_para.font.name = title_para1.font.name = "Arial"
    title_para.font.bold = title_para1.font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = title.text_frame.paragraphs[1].font.color.rgb = maincolor

    #bottom right text
    companybox_title = slide.shapes.add_textbox(left=Cm(20),top=Cm(17.5), width=Cm(10), height=Cm(1))
    tf = companybox_title.text_frame
    tf.text = bottomrighttitle.replace("\n"," ")
    tf.paragraphs[0].font.color.rgb = maincolor
    tf.paragraphs[0].font.size = Pt(11)

    try:
        logo = slide.shapes.add_picture("logo.png", left=Inches(1), top=Cm(2), height=Cm(4))
        logo.left=int(round((pres.slide_width - logo.width) / 2))
        print("logo.png found")
    except FileNotFoundError:
        print("logo.png not found")


    #all other slides
    index = 1
    header_slide_layout = pres.slide_layouts[5]#Header slide layout
    for i in information: #for product dictionary in the dictionary list
        #Slide and title
        slide = pres.slides.add_slide(header_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = i["title"]
        title_shape_para = slide.shapes.title.text_frame.paragraphs[0]
        title_shape_para.font.name = "Arial"
        title_shape_para.font.size = Pt(40)
        title_shape.text_frame.paragraphs[0].font.color.rgb = maincolor

        #Logo top right per slide, sets it as background
        try:
            logo_small = slide.shapes.add_picture("logo.png", left=Cm(21), top=Cm(1), height=Cm(3))
            slide.shapes._spTree.remove(logo_small._element)
            slide.shapes._spTree.insert(2, logo_small._element)
        except FileNotFoundError:
            pass

        #Image per slide
        left = Inches(1.25)
        top = Inches(3)
        pic = slide.shapes.add_picture("AmazonProductImage"+str(index)+".jpg", left, top, height=Inches(2.5))

        #Price per slide
        pricebox = slide.shapes.add_textbox(left=left, top=Inches(2.25), width=Cm(10), height=Inches(2.5))
        price = pricebox.text_frame
        price.text = i["price"]
        price.paragraphs[0].font.color.rgb = maincolor
        price.paragraphs[0].font.size = Pt(28)

        """
        #Url under image
        urlbox = slide.shapes.add_textbox(left=left,top=Cm(14), width=Cm(10), height=Cm(2))
        tf = urlbox.text_frame
        tf.text = i["url"]
        """

        #Specifications text box
        specificationsbox = slide.shapes.add_textbox(left=Cm(13), top=Inches(2.25), width=Cm(10), height=Inches(2.5))
        specifications = specificationsbox.text_frame
        specifications.text = "Product specifications"
        specifications.paragraphs[0].font.color.rgb = maincolor
        specifications.paragraphs[0].font.size = Pt(26)
        number_of_para = 1
        for kvpair in i["characteristics"]:#for specification in specifications dictionary
            p = specifications.add_paragraph()
            stringedspecs = str(kvpair).replace('{','').replace('}','').replace('\'','')
            
            characterinline = 0
            #if a line in the string is longer than 38 characters, the last space becomes "\n". This way the text won't go outside of the slide
            if len(stringedspecs) > 38:
                for c in stringedspecs:
                    if c == " ":
                        characterinline = stringedspecs.index(c)
                stringedspecs = stringedspecs[:characterinline]+"-\n"+stringedspecs[characterinline+1:]

            p.text = stringedspecs
            specifications.paragraphs[number_of_para].font.size = Pt(20)
            specifications.paragraphs[number_of_para].font.color.rgb = secondcolor
            number_of_para += 1

        #Bottom right text every slide
        companybox = slide.shapes.add_textbox(left=Cm(20),top=Cm(17.5), width=Cm(10), height=Cm(1))
        tf = companybox.text_frame
        tf.text = bottomrighttitle.replace("\n"," ")
        tf.paragraphs[0].font.color.rgb = maincolor
        tf.paragraphs[0].font.size = Pt(11)

        index += 1
    try:
        pres.save(f"{filename}.pptx")
    except PermissionError:
        sys.exit("Couldn't save, due to the Powerpoint file still being open. Please close the Powerpoint file and try again")

def imageremover(productnumber):
    for n in range(1, productnumber+1):
        os.remove(f"AmazonProductImage{n}.jpg")
    
    imagesremoved = True
    for n in range(1, 1000):
        b = os.path.exists(f"AmazonProductImage{n}.jpg")#If there are still AmazonProductImages
        if b:
            imagesremoved = False
    return imagesremoved

if __name__ == "__main__":
    main() 